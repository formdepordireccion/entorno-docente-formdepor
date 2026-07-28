#!/usr/bin/env python3
"""Genera un .pptx a partir de un Markdown con secciones `##`, aplicando el
sistema de diseño "UD00 style" (naranja/Arial) descrito en
`.claude/skills/asig-recursos/reference/formdepor-pptx-ud00-style.md`.

Uso:
    .venv-recursos/bin/python .claude/skills/asig-recursos/scripts/generar_pptx.py \
        --input RUTA_MARKDOWN --output RUTA_PPTX --titulo "Título de la presentación"

No hay límite de diapositivas: prioriza cubrir el temario real con
detalle sobre mantener el esquema corto — es preferible partir un bloque
denso en 2-3 diapositivas antes que forzarlo en una sola.

Regla de conversión del esquema (deliberadamente simple, sin red):
- El primer `#` del archivo (si existe) se ignora como título de sección;
  el título real de la portada lo da siempre --titulo.
- Cada `##` es una diapositiva nueva. Tipos de sección, detectados por el
  propio título:
  - **Separador de bloque** (Patrón D): el título empieza por `RA<número>`
    (p. ej. `## RA7 — Valoración inicial`, extrae "RA7" como
    macroidentificador) o por la palabra `Separador` (para secciones sin
    RA asociado). Fondo oscuro, sin viñetas.
  - **Cierre** (Patrón N): el título empieza por `Cierre`. Fondo oscuro,
    tarjeta naranja de título (+ hasta 2 líneas de subtítulo, ver abajo)
    y hasta 3 próximos pasos con check.
  - **Protocolo de tres pasos** (Patrón F, automático): cualquier otra
    sección con EXACTAMENTE 3 viñetas, todas en formato `LETRA/Nº: texto`
    (p. ej. `- P: Proteger — tu propia seguridad primero`). Colores de
    círculo rotan naranja/rojo/negro.
  - **Contenido con patrón elegido a mano**: quien redacta el esquema
    (paso 1 del proceso, en `SKILL.md`) decide, diapositiva a diapositiva,
    qué patrón representa mejor ese contenido — igual que pide la sección
    13 de la guía ("seleccionar para cada diapositiva uno de los patrones
    A-N") — añadiendo una etiqueta al final del título de la sección:
    - `## Título {recorrido}` — flechas tipo chevron naranja→negro→gris
      (Patrón C). Para recorridos de 3-5 etapas (mapa de la unidad,
      cadena de supervivencia...). Usa las dos primeras líneas `>` como
      frase de contexto + subrayado, con una etiqueta "RECORRIDO" encima
      de las flechas.
    - `## Título {clasificacion}` — columnas con cabecera de color,
      cada una con su categoría (Patrón G/K). Para taxonomías (tipos de
      lesión, grupos del botiquín, roles musculares...).
    - `## Título {comparativa}` — dos paneles grandes enfrentados, uno
      oscuro y uno naranja (Patrón H). Para exactamente 2 conceptos que
      se contraponen (agudo/crónico, antes/después...). Si el cuerpo de
      una viñeta tiene varios datos separados por `;`, se listan como
      mini-viñetas dentro del panel.
    - `## Título {reglaexcepciones}` — banner rojo con la regla general
      (primera viñeta) + tarjetas numeradas con las excepciones (resto
      de viñetas) (Patrón L). Para "no hagas X salvo que...".
    - Sin etiqueta → **tarjetas** (por defecto): cuadrícula de tarjetas
      con círculo numerado, la más versátil para contenido genérico.
    Máximo 5 viñetas de contenido por diapositiva en cualquier patrón —
    si hay más, se trocea automáticamente en varias diapositivas
    `(1/2)`, `(2/2)`... (el subtítulo y el callout de la sección viajan
    siempre con la primera/última, nunca solos en una diapositiva aparte).
  - **Subtítulo de diapositiva**: una línea `> texto` justo después del
    título de sección (antes de las viñetas) se muestra como frase de
    contexto bajo la línea de acento naranja. `{recorrido}` admite dos
    líneas `>` (la primera en negrita, la segunda más tenue); `{cierre}`
    admite hasta dos como cuerpo de la tarjeta naranja; el resto usa solo
    la primera.
  - **Caja de cierre de diapositiva**: una viñeta `- ! texto` (al final
    de la sección) no se dibuja como tarjeta — se extrae y se muestra
    como una caja horizontal a pie de diapositiva. Si el texto es
    `Etiqueta: frase` (mismo convenio que las tarjetas), sale en caja gris
    clara con la etiqueta en naranja ("IDEA CLAVE", "OJO"...); si no,
    en barra negra sólida centrada.
- Cada diapositiva de contenido lleva un icono arriba a la derecha,
  elegido por palabra clave del título (ver `ICONOS_POR_PALABRA_CLAVE`)
  entre un set de formas nativas de PowerPoint (sección 7.6 de la guía:
  "estilo lineal, trazo simple" — no son un set de iconos real, es una
  aproximación razonable sin depender de assets externos).
- Las líneas que empiezan por `-`/`*` son viñetas; los `###` se añaden
  como subtítulo en negrita sin crear diapositiva nueva. Emphasis
  Markdown (`**negrita**`, `*cursiva*`, `` `code` ``) se limpia del texto.

Patrones del documento de diseño que este script **no** implementa
(quedarían para maquetación manual en PowerPoint, o para una futura
plantilla maestra real): B, E, I, M — requieren estructura (2×2 con
cuatro bloques temáticos fijos, bifurcación de decisión con dos
resultados coloreados...) que no tiene un equivalente natural en un
esquema lineal de título+viñetas. Las fotografías reales tampoco se
generan — eso se añade a mano después.
"""

import argparse
import re
import unicodedata

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --- Paleta "UD00 style" (sección 3.1 de la guía) --------------------------

BRAND_ORANGE = RGBColor(0xFF, 0x6A, 0x00)
INK_BLACK = RGBColor(0x14, 0x16, 0x1A)
INK_RAISED = RGBColor(0x1E, 0x22, 0x29)
PAPER_LIGHT = RGBColor(0xF2, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x2B, 0x2F, 0x36)
SLATE = RGBColor(0x4A, 0x50, 0x58)
MUTED_GRAY = RGBColor(0x6C, 0x75, 0x7D)
SOFT_GRAY = RGBColor(0xB9, 0xBE, 0xC6)
LINE_GRAY = RGBColor(0xD8, 0xDB, 0xDF)

FUENTE = "Arial"

MAX_PALABRAS_TITULO_PORTADA = 12
MAX_VINETAS_POR_DIAPOSITIVA = 5

avisos = []


def limpiar_markdown(texto: str) -> str:
    texto = re.sub(r"\*\*(.+?)\*\*", r"\1", texto)
    texto = re.sub(r"\*(.+?)\*", r"\1", texto)
    texto = re.sub(r"`(.+?)`", r"\1", texto)
    return texto.strip()


def normalizar(texto: str) -> str:
    sin_tildes = unicodedata.normalize("NFKD", texto).encode("ascii", "ignore").decode("ascii")
    return sin_tildes.lower().strip()


def quitar_prefijo_separador(titulo: str) -> str:
    return re.sub(r"^separador\s*[:\-–—]*\s*", "", titulo, flags=re.IGNORECASE).strip()


def extraer_macroidentificador(titulo: str):
    """'RA7 — Valoración inicial' -> ('RA7', 'Valoración inicial')."""
    m = re.match(r"^(RA\d+)\s*[:\-–—]*\s*(.*)$", titulo, flags=re.IGNORECASE)
    if m:
        return m.group(1).upper(), m.group(2).strip()
    return None, titulo


def parsear_secciones(md_texto: str):
    secciones = []
    titulo_actual = None
    contenido_actual = []
    for linea in md_texto.splitlines():
        m2 = re.match(r"^##\s+(.*)", linea)
        if m2:
            if titulo_actual is not None:
                secciones.append((titulo_actual, contenido_actual))
            titulo_actual = limpiar_markdown(m2.group(1))
            contenido_actual = []
            continue
        if titulo_actual is None:
            continue
        m3 = re.match(r"^###\s+(.*)", linea)
        if m3:
            contenido_actual.append(("sub", limpiar_markdown(m3.group(1))))
            continue
        mq = re.match(r"^>\s+(.*)", linea)
        if mq:
            contenido_actual.append(("subtitulo", limpiar_markdown(mq.group(1))))
            continue
        mb = re.match(r"^\s*[-*]\s+(.*)", linea)
        if mb:
            texto_bullet = limpiar_markdown(mb.group(1))
            mc = re.match(r"^!\s*(.*)", texto_bullet)
            if mc:
                contenido_actual.append(("callout", mc.group(1)))
            else:
                contenido_actual.append(("bullet", texto_bullet))
            continue
        if linea.strip():
            contenido_actual.append(("bullet", limpiar_markdown(linea.strip())))
    if titulo_actual is not None:
        secciones.append((titulo_actual, contenido_actual))
    return secciones


def separar_bloque(bloque):
    """([subtítulos en orden], callout|None, [(tipo, texto) solo 'bullet'/'sub'])."""
    subtitulos = [t for tipo, t in bloque if tipo == "subtitulo"]
    callout = next((t for tipo, t in bloque if tipo == "callout"), None)
    items = [(tipo, t) for tipo, t in bloque if tipo in ("bullet", "sub")]
    return subtitulos, callout, items


def anadir_subtitulo(slide, texto, color=SLATE):
    if not texto:
        return
    anadir_texto(slide, texto, Inches(0.75), Inches(1.32), Inches(10.2), Inches(0.4), 15, color)


CALLOUT_CLARO = RGBColor(0xE9, 0xEA, 0xEC)


def anadir_callout(slide, texto, top):
    """Caja de cierre de diapositiva (sección 'Idea clave'/'Ojo' de la
    referencia). Si el texto viene como 'Etiqueta: frase' (mismo convenio
    que `dividir_titulo_cuerpo`), se muestra en caja gris clara con la
    etiqueta en naranja; si no, en barra negra sólida centrada."""
    if not texto:
        return top
    etiqueta, cuerpo = dividir_titulo_cuerpo(texto)
    alto = Inches(0.55) if len(texto) < 85 else Inches(0.78)
    caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), top, Inches(11.83), alto)
    caja.fill.solid()
    caja.line.fill.background()
    caja.shadow.inherit = False
    tf = caja.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]

    if etiqueta:
        caja.fill.fore_color.rgb = CALLOUT_CLARO
        p.alignment = PP_ALIGN.LEFT
        run_etiqueta = p.add_run()
        run_etiqueta.text = f"{etiqueta.upper()}   "
        estilo_run(run_etiqueta, 13, BRAND_ORANGE, negrita=True)
        run_texto = p.add_run()
        run_texto.text = cuerpo
        estilo_run(run_texto, 13, CHARCOAL, negrita=False)
    else:
        caja.fill.fore_color.rgb = INK_BLACK
        p.alignment = PP_ALIGN.CENTER
        run_texto = p.add_run()
        run_texto.text = texto
        estilo_run(run_texto, 13, WHITE, negrita=True)

    return top + alto + Inches(0.12)


def cabecera_contenido(slide, titulo, subtitulo, color_titulo=CHARCOAL, color_icono=BRAND_ORANGE):
    """Título + línea de acento + subtítulo opcional + icono — cabecera
    común a todas las diapositivas de contenido claro. Devuelve la
    coordenada Y donde debe empezar el cuerpo."""
    anadir_texto(slide, titulo, Inches(0.5), Inches(0.3), Inches(10.6), Inches(0.75), 32, color_titulo, negrita=True)
    anadir_linea(slide, BRAND_ORANGE, Inches(0.75), Inches(1.08), Inches(0.8), Pt(4))
    anadir_icono(slide, titulo, Inches(11.9), Inches(0.35), Inches(0.55), color_icono)
    if subtitulo:
        anadir_texto(slide, subtitulo, Inches(0.75), Inches(1.26), Inches(9.9), Inches(0.4), 15, SLATE)
        return Inches(1.78)
    return Inches(1.45)


def extraer_patron(titulo: str):
    """'Niveles de prevención {proceso}' -> ('Niveles de prevención', 'proceso')."""
    m = re.match(r"^(.*?)\s*\{(\w+)\}\s*$", titulo)
    if m:
        return m.group(1).strip(), m.group(2).lower()
    return titulo, None


def clasificar_seccion(titulo, items):
    titulo_normalizado = normalizar(titulo)
    if titulo_normalizado.startswith("cierre"):
        return "cierre"
    if titulo_normalizado.startswith("separador") or re.match(r"^ra\d+\b", titulo_normalizado):
        return "separador"
    bullets = [t for tipo, t in items if tipo == "bullet"]
    if len(bullets) == 3:
        patrones = [re.match(r"^([A-ZÁÉÍÓÚÑ0-9]{1,3})\s*[:\-–—]\s*(.+)$", b) for b in bullets]
        if all(patrones):
            return "protocolo"
    return "contenido"


# --- Helpers de estilo ------------------------------------------------------

def estilo_run(run, tamano_pt, color, negrita=None, mayusculas=False):
    run.font.name = FUENTE
    run.font.size = Pt(tamano_pt)
    run.font.color.rgb = color
    if negrita is not None:
        run.font.bold = negrita
    if mayusculas:
        run.text = run.text.upper()


def fondo_solido(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def anadir_linea(slide, color, left, top, width, height=Pt(3)):
    forma = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    forma.fill.solid()
    forma.fill.fore_color.rgb = color
    forma.line.fill.background()
    forma.shadow.inherit = False
    return forma


def anadir_texto(slide, texto, left, top, width, height, tamano_pt, color, negrita=False,
                  alineacion=PP_ALIGN.LEFT, mayusculas=False):
    caja = slide.shapes.add_textbox(left, top, width, height)
    tf = caja.text_frame
    tf.word_wrap = True
    tf.text = texto.upper() if mayusculas else texto
    p = tf.paragraphs[0]
    p.alignment = alineacion
    estilo_run(p.runs[0], tamano_pt, color, negrita=negrita)
    return caja


def anadir_pie(slide, prs, pie_izquierda, numero, color_texto):
    top = prs.slide_height - Inches(0.38)
    anadir_texto(slide, pie_izquierda, Inches(0.6), top, Inches(7), Inches(0.3),
                 9, color_texto, negrita=True, mayusculas=True)
    caja = slide.shapes.add_textbox(prs.slide_width - Inches(1.5), top, Inches(0.9), Inches(0.3))
    tf = caja.text_frame
    tf.text = str(numero)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    estilo_run(tf.paragraphs[0].runs[0], 10, BRAND_ORANGE, negrita=True)


def revisar_titulo_portada(titulo):
    n = len(titulo.split())
    if n > MAX_PALABRAS_TITULO_PORTADA:
        avisos.append(
            f"Título de portada \"{titulo}\" tiene {n} palabras (máximo recomendado: {MAX_PALABRAS_TITULO_PORTADA}, sección 8/Patrón A de la guía)."
        )


# --- Iconografía (sección 7.6 de la guía: "estilo lineal, trazo simple") ----
#
# No hay assets de icono reales disponibles (SVG/PNG) — se construyen con
# formas nativas de PowerPoint combinadas, en línea (sin relleno), para
# mantener el espíritu "estilo lineal, trazo simple y uniforme" sin
# depender de una librería de iconos externa. No son réplica de un set de
# iconos concreto, son una aproximación razonable por temática.

def _forma_linea(slide, tipo, x, y, size, color, grosor=Pt(2.25), rotacion=0, relleno=False):
    forma = slide.shapes.add_shape(tipo, x, y, size, size)
    if relleno:
        forma.fill.solid()
        forma.fill.fore_color.rgb = color
        forma.line.fill.background()
    else:
        forma.fill.background()
        forma.line.color.rgb = color
        forma.line.width = grosor
    forma.rotation = rotacion
    forma.shadow.inherit = False
    return forma


def icono_corazon(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.HEART, x, y, size, color, relleno=relleno)


def icono_cruz(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.CROSS, x, y, size, color, relleno=relleno)


def icono_alerta(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, size, color, relleno=relleno)
    color_texto = WHITE if relleno else color
    caja = slide.shapes.add_textbox(x, y + size * 0.28, size, size * 0.5)
    caja.text_frame.text = "!"
    caja.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    estilo_run(caja.text_frame.paragraphs[0].runs[0], int(size / 152400 * 11), color_texto, negrita=True)


def icono_prohibido(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.NO_SYMBOL, x, y, size, color, relleno=relleno)


def icono_rayo(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.LIGHTNING_BOLT, x, y, size, color, relleno=relleno)


def icono_escudo(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.PENTAGON, x, y, size, color, rotacion=90, relleno=relleno)


def icono_reloj(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.OVAL, x, y, size, color)
    cx, cy = x + int(size / 2), y + int(size / 2)
    manecilla1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy, cx, y + int(size * 0.22))
    manecilla1.line.color.rgb = color
    manecilla1.line.width = Pt(2)
    manecilla2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy, x + int(size * 0.72), cy)
    manecilla2.line.color.rgb = color
    manecilla2.line.width = Pt(2)


def icono_circular(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.CIRCULAR_ARROW, x, y, size, color, relleno=relleno)


def icono_nube(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.CLOUD, x, y, size, color, relleno=relleno)


def icono_hueso(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.OVAL, x, y, int(size * 0.34), color, grosor=Pt(1.75))
    _forma_linea(slide, MSO_SHAPE.OVAL, x + int(size * 0.66), y + int(size * 0.66), int(size * 0.34), color, grosor=Pt(1.75))
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + int(size * 0.22), y + int(size * 0.22), int(size * 0.56), Pt(3))
    barra.rotation = 45
    barra.fill.solid()
    barra.fill.fore_color.rgb = color
    barra.line.fill.background()
    barra.shadow.inherit = False


def icono_calendario(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, color, grosor=Pt(2))
    for i in (0.28, 0.62):
        anilla = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + int(size * i), y - Pt(2), Pt(3), Pt(10))
        anilla.fill.solid()
        anilla.fill.fore_color.rgb = color
        anilla.line.fill.background()
        anilla.shadow.inherit = False


def icono_estrella(slide, x, y, size, color, relleno=False):
    _forma_linea(slide, MSO_SHAPE.STAR_5_POINT, x, y, size, color, relleno=relleno)


ICONOS_POR_PALABRA_CLAVE = [
    (("cadena pas", "avisar", "112", "aviso"), icono_rayo),
    (("botiquin", "medicamento", "primeros auxilios", "curacion", "antiseptic"), icono_cruz),
    (("legal", "etico", "derecho", "obligacion", "psicologic"), icono_escudo),
    (("cardio", "corazon", "circulator", "pulso", "gasto cardiaco"), icono_corazon),
    (("lesion", "herida", "riesgo", "quemadura", "hemorrag", "trastorno"), icono_alerta),
    (("movilizar", "movilizacion", "inmoviliza", "no mover"), icono_prohibido),
    (("valoracion", "sistematica", "diagnostic", "reconoc"), icono_reloj),
    (("cierre", "proximos pasos", "bienvenida"), icono_calendario),
    (("hueso", "oseo", "esqueletic"), icono_hueso),
    (("articula", "movimiento", "palanca", "biomecanic", "grado de libertad", "plano"), icono_circular),
    (("respirat", "pulmon", "via", "alveolo"), icono_nube),
    (("metabolismo", "energia", "fisiolog", "muscul", "fibra", "contraccion"), icono_rayo),
]


def icono_para(texto):
    t = normalizar(texto)
    for claves, builder in ICONOS_POR_PALABRA_CLAVE:
        if any(clave in t for clave in claves):
            return builder
    return icono_estrella


def anadir_icono(slide, texto_contexto, x, y, size, color, relleno=False):
    icono_para(texto_contexto)(slide, x, y, size, color, relleno=relleno)


# --- Diapositivas ------------------------------------------------------------

def slide_portada(prs, titulo, subtitulo, etiqueta, pie_izquierda, numero):
    revisar_titulo_portada(titulo)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, INK_BLACK)

    if etiqueta:
        anadir_texto(slide, etiqueta, Inches(0.75), Inches(1.3), Inches(8), Inches(0.4),
                     13, BRAND_ORANGE, negrita=True, mayusculas=True)
    anadir_linea(slide, BRAND_ORANGE, Inches(0.75), Inches(1.85), Inches(1.1), Pt(4))
    anadir_texto(slide, titulo, Inches(0.75), Inches(2.1), Inches(8.3), Inches(2.2),
                 40, WHITE, negrita=True)
    if subtitulo:
        anadir_texto(slide, subtitulo, Inches(0.75), Inches(4.6), Inches(8), Inches(0.6),
                     18, BRAND_ORANGE, negrita=False)
    anadir_linea(slide, BRAND_ORANGE, Inches(0.75), Inches(5.55), Inches(4.5), Pt(4))

    tarjeta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.9), Inches(1.75), Inches(2.7), Inches(4))
    tarjeta.fill.solid()
    tarjeta.fill.fore_color.rgb = INK_RAISED
    tarjeta.line.fill.background()
    tarjeta.shadow.inherit = False
    anadir_icono(slide, titulo + " " + (etiqueta or ""), Inches(10.65), Inches(3.15), Inches(1.2), BRAND_ORANGE, relleno=True)
    anadir_linea(slide, BRAND_ORANGE, Inches(9.9), Inches(1.55), Inches(2.7), Pt(4))

    anadir_pie(slide, prs, pie_izquierda, numero, SOFT_GRAY)
    return slide


def slide_separador(prs, titulo, pie_izquierda, numero, bloque_actual, bloque_total):
    macro_id, titulo_seccion = extraer_macroidentificador(quitar_prefijo_separador(titulo))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, INK_BLACK)
    anadir_texto(slide, f"BLOQUE {bloque_actual} DE {bloque_total}", Inches(0.75), Inches(2.05),
                 Inches(5.5), Inches(0.35), 13, BRAND_ORANGE, negrita=True, mayusculas=True)
    if macro_id:
        anadir_texto(slide, macro_id, Inches(0.75), Inches(2.55), Inches(4.44), Inches(1.7),
                     90, BRAND_ORANGE, negrita=True)
    anadir_linea(slide, BRAND_ORANGE, Inches(5.6), Inches(2.3), Pt(3), Inches(2.2))
    anadir_texto(slide, titulo_seccion, Inches(5.97), Inches(2.7), Inches(6.61), Inches(1.8),
                 30, WHITE, negrita=True)
    anadir_linea(slide, BRAND_ORANGE, Inches(0.75), Inches(6.9), Inches(4.0), Pt(3))
    anadir_pie(slide, prs, pie_izquierda, numero, SOFT_GRAY)
    return slide


def slide_cierre(prs, titulo, bloque, pie_izquierda, numero):
    subtitulos, _callout, items = separar_bloque(bloque)
    acciones = [t for tipo, t in items if tipo == "bullet"][:3]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, INK_BLACK)

    tarjeta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.2),
                                      Inches(5.3), Inches(5.2))
    tarjeta.fill.solid()
    tarjeta.fill.fore_color.rgb = BRAND_ORANGE
    tarjeta.line.fill.background()
    tarjeta.shadow.inherit = False
    tf = tarjeta.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.35)
    tf.text = titulo
    estilo_run(tf.paragraphs[0].runs[0], 28, WHITE, negrita=True)
    for texto_sub in subtitulos[:2]:
        p = tf.add_paragraph()
        p.text = texto_sub
        p.space_before = Pt(10)
        estilo_run(p.runs[0], 14, WHITE, negrita=False)

    top = Inches(1.4)
    for texto in acciones:
        marca = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), top, Inches(0.34), Inches(0.34))
        marca.fill.solid()
        marca.fill.fore_color.rgb = BRAND_ORANGE
        marca.line.fill.background()
        marca.shadow.inherit = False
        tfm = marca.text_frame
        tfm.text = "✓"
        tfm.paragraphs[0].alignment = PP_ALIGN.CENTER
        estilo_run(tfm.paragraphs[0].runs[0], 15, WHITE, negrita=True)
        anadir_texto(slide, texto, Inches(7.2), top - Inches(0.07), Inches(5.35), Inches(0.9), 16, WHITE)
        top += Inches(1.05)

    anadir_pie(slide, prs, pie_izquierda, numero, SOFT_GRAY)
    return slide


PALETA_PROTOCOLO = [BRAND_ORANGE, RGBColor(0xE6, 0x39, 0x46), INK_BLACK]


def slide_protocolo(prs, titulo, bloque, pie_izquierda, numero):
    """Patrón F — protocolo de tres pasos, tarjeta+círculo con color rotado
    naranja/rojo/negro (igual que la referencia UD00, sección PAS)."""
    subtitulos, callout, items = separar_bloque(bloque)
    bullets = [t for tipo, t in items if tipo == "bullet"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, PAPER_LIGHT)
    y0 = cabecera_contenido(slide, titulo, subtitulos[0] if subtitulos else None)

    columnas_x = [Inches(0.75), Inches(4.68), Inches(8.61)]
    col_w = Inches(3.4)
    card_h = Inches(3.35)
    parseados = []
    for texto in bullets:
        m = re.match(r"^([A-ZÁÉÍÓÚÑ0-9]{1,3})\s*[:\-–—]\s*(.+)$", texto)
        parseados.append((m.group(1), m.group(2)) if m else (texto[:1], texto))

    for i, (letra, explicacion) in enumerate(parseados):
        x = columnas_x[i]
        color = PALETA_PROTOCOLO[i % len(PALETA_PROTOCOLO)]

        tarjeta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, col_w, card_h)
        tarjeta.fill.solid()
        tarjeta.fill.fore_color.rgb = WHITE
        tarjeta.line.color.rgb = LINE_GRAY
        tarjeta.line.width = Pt(0.75)
        tarjeta.shadow.inherit = False
        tarjeta.text_frame.text = ""

        circulo = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y0 + Inches(0.3), Inches(0.9), Inches(0.9))
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = color
        circulo.line.fill.background()
        circulo.shadow.inherit = False
        tf = circulo.text_frame
        tf.text = letra
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        estilo_run(tf.paragraphs[0].runs[0], 28, WHITE, negrita=True)

        accion, _, descripcion = explicacion.partition(" — ")
        if not descripcion:
            accion, _, descripcion = explicacion.partition(" - ")
        if not descripcion:
            accion, descripcion = "", explicacion

        anadir_texto(slide, accion.upper() if accion else letra, x + Inches(0.28), y0 + Inches(1.4),
                     col_w - Inches(0.55), Inches(0.4), 17, color, negrita=True)
        anadir_texto(slide, descripcion, x + Inches(0.28), y0 + Inches(1.85),
                     col_w - Inches(0.55), card_h - Inches(2.1), 13, CHARCOAL)

    bottom = y0 + card_h + Inches(0.25)
    if callout:
        anadir_callout(slide, callout, bottom)
    anadir_pie(slide, prs, pie_izquierda, numero, MUTED_GRAY)
    return slide


PALETA_RECORRIDO = [BRAND_ORANGE, INK_BLACK, MUTED_GRAY, RGBColor(0x00, 0xA8, 0xB5), RGBColor(0xE6, 0x39, 0x46)]


def slide_recorrido(prs, titulo, bloque, pie_izquierda, numero):
    """Patrón C — recorrido con flechas tipo chevron, naranja→negro→gris
    (Pattern C de la guía: "naranja para el bloque activo o inicial;
    negro para el bloque central; gris para el bloque final")."""
    subtitulos, callout, items = separar_bloque(bloque)
    pasos = [t for tipo, t in items if tipo == "bullet"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, PAPER_LIGHT)

    anadir_texto(slide, titulo, Inches(0.5), Inches(0.3), Inches(10.6), Inches(0.55), 28, CHARCOAL, negrita=True)
    top_texto = Inches(1.0)
    for texto_sub in subtitulos[:2]:
        negrita = texto_sub == subtitulos[0]
        color = CHARCOAL if negrita else SLATE
        anadir_texto(slide, texto_sub, Inches(0.75), top_texto, Inches(10.8), Inches(0.4), 16, color, negrita=negrita)
        top_texto += Inches(0.42)

    y_eyebrow = top_texto + Inches(0.25)
    anadir_texto(slide, "RECORRIDO", Inches(0.75), y_eyebrow, Inches(6), Inches(0.35),
                 12, BRAND_ORANGE, negrita=True, mayusculas=True)

    n = len(pasos)
    content_x0, content_w = Inches(0.75), Inches(11.83)
    gap = Inches(0.06)
    col_w = int((content_w - (n - 1) * gap) / n)
    y0 = y_eyebrow + Inches(0.5)
    alto = Inches(2.2)

    for i, texto in enumerate(pasos):
        x = content_x0 + i * (col_w + gap)
        color = PALETA_RECORRIDO[i % len(PALETA_RECORRIDO)]
        titulo_paso, cuerpo_paso = dividir_titulo_cuerpo(texto)

        chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y0, col_w, alto)
        chevron.fill.solid()
        chevron.fill.fore_color.rgb = color
        chevron.line.fill.background()
        chevron.shadow.inherit = False
        chevron.text_frame.text = ""

        # El texto va en una caja aparte, superpuesta: la zona de texto
        # "segura" que PowerPoint calcula para un CHEVRON se basa en la
        # geometría de la flecha (muy estrecha por la punta), no en el
        # ancho real de la forma — poner el texto directamente en
        # chevron.text_frame provoca un ajuste de línea catastrófico
        # (una letra por línea). Con una caja propia controlamos el
        # ancho nosotros.
        texto_w = int(col_w * 0.62)
        texto_x = x + int((col_w - texto_w) / 2)
        caja = slide.shapes.add_textbox(texto_x, y0, texto_w, alto)
        tf = caja.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = (titulo_paso or cuerpo_paso).upper()
        estilo_run(p1.runs[0], 17, WHITE, negrita=True)
        if titulo_paso and cuerpo_paso:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = cuerpo_paso
            p2.space_before = Pt(4)
            estilo_run(p2.runs[0], 10.5, WHITE)

    bottom = y0 + alto + Inches(0.3)
    if callout:
        anadir_callout(slide, callout, bottom)
    anadir_pie(slide, prs, pie_izquierda, numero, MUTED_GRAY)
    return slide


def slide_reglaexcepciones(prs, titulo, bloque, pie_izquierda, numero):
    """Patrón L — banner con la regla general + tarjetas numeradas con las
    excepciones. La primera viñeta es la regla; el resto, excepciones."""
    subtitulos, callout, items = separar_bloque(bloque)
    bullets = [t for tipo, t in items if tipo == "bullet"]
    if not bullets:
        return slide_contenido(prs, titulo, bloque, pie_izquierda, numero)
    regla, excepciones = bullets[0], bullets[1:]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, PAPER_LIGHT)
    y0 = cabecera_contenido(slide, titulo, subtitulos[0] if subtitulos else None)

    etiqueta_regla, texto_regla = dividir_titulo_cuerpo(regla)
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), y0, Inches(11.83), Inches(1.15))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0xE6, 0x39, 0x46)
    banner.line.fill.background()
    banner.shadow.inherit = False
    tf = banner.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    if etiqueta_regla:
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.text = etiqueta_regla.upper()
        estilo_run(p0.runs[0], 12, RGBColor(0xFF, 0xD5, 0xD9), negrita=True)
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.text = texto_regla
        estilo_run(p1.runs[0], 20, WHITE, negrita=True)
    else:
        tf.text = texto_regla
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        estilo_run(tf.paragraphs[0].runs[0], 20, WHITE, negrita=True)

    n = len(excepciones)
    y_exc = y0 + Inches(1.45)
    if n:
        anadir_texto(slide, "LAS EXCEPCIONES", Inches(0.75), y_exc, Inches(6), Inches(0.35),
                     12, BRAND_ORANGE, negrita=True, mayusculas=True)
        y_exc += Inches(0.45)
        content_w = Inches(11.83)
        gap = Inches(0.22)
        card_w = int((content_w - (n - 1) * gap) / n)
        card_h = Inches(1.75)
        for i, texto in enumerate(excepciones):
            x = Inches(0.75) + i * (card_w + gap)
            titulo_exc, cuerpo_exc = dividir_titulo_cuerpo(texto)

            tarjeta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_exc, card_w, card_h)
            tarjeta.fill.solid()
            tarjeta.fill.fore_color.rgb = WHITE
            tarjeta.line.color.rgb = LINE_GRAY
            tarjeta.line.width = Pt(0.75)
            tarjeta.shadow.inherit = False
            tarjeta.text_frame.text = ""

            circulo = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y_exc + Inches(0.18), Inches(0.4), Inches(0.4))
            circulo.fill.solid()
            circulo.fill.fore_color.rgb = BRAND_ORANGE
            circulo.line.fill.background()
            circulo.shadow.inherit = False
            tfc = circulo.text_frame
            tfc.text = str(i + 1)
            tfc.paragraphs[0].alignment = PP_ALIGN.CENTER
            estilo_run(tfc.paragraphs[0].runs[0], 14, WHITE, negrita=True)

            anadir_texto(slide, titulo_exc or cuerpo_exc, x + Inches(0.7), y_exc + Inches(0.16),
                         card_w - Inches(0.9), Inches(0.4), 14, CHARCOAL, negrita=True)
            if titulo_exc:
                anadir_texto(slide, cuerpo_exc, x + Inches(0.22), y_exc + Inches(0.68),
                             card_w - Inches(0.44), card_h - Inches(0.8), 11.5, SLATE)
        y_exc += card_h

    if callout:
        anadir_callout(slide, callout, y_exc + Inches(0.25))
    anadir_pie(slide, prs, pie_izquierda, numero, MUTED_GRAY)
    return slide


def dividir_titulo_cuerpo(texto):
    """'Primaria: evitar que el accidente ocurra' -> ('Primaria', 'evitar que el accidente ocurra').

    Si la viñeta no tiene 'concepto: explicación', se trata como cuerpo sin
    micro-título (la tarjeta solo lleva el número y el texto).
    """
    m = re.match(r"^([^:]{2,28}):\s*(.+)$", texto)
    if m:
        return m.group(1).strip(), m.group(2).strip()
    return None, texto


def filas_de_tarjetas(n):
    """Reparte n tarjetas (1-5, límite de MAX_VINETAS_POR_DIAPOSITIVA) en
    filas de máx. 3 columnas."""
    return {1: [1], 2: [2], 3: [3], 4: [2, 2], 5: [3, 2]}[n]


def anadir_lista(slide, items, left, top, width, height, tamano_pt, color):
    """Textbox con un párrafo por item, cada uno con viñeta '—'."""
    caja = slide.shapes.add_textbox(left, top, width, height)
    tf = caja.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"—  {item}"
        p.space_after = Pt(6)
        estilo_run(p.runs[0], tamano_pt, color)
    return caja


PALETA_CLASIFICACION = [BRAND_ORANGE, INK_BLACK, RGBColor(0x00, 0xA8, 0xB5), RGBColor(0xE6, 0x39, 0x46), MUTED_GRAY]


def slide_clasificacion(prs, titulo, bloque, pie_izquierda, numero):
    """Patrón G/K — columnas con cabecera de color, una por categoría."""
    subtitulos, callout, items_raw = separar_bloque(bloque)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, PAPER_LIGHT)
    y0 = cabecera_contenido(slide, titulo, subtitulos[0] if subtitulos else None)

    items = [texto for tipo, texto in items_raw if tipo == "bullet"]
    n = len(items)
    gap = Inches(0.22)
    content_x0, content_w = Inches(0.75), Inches(11.83)
    col_w = int((content_w - (n - 1) * gap) / n)
    header_h = Inches(0.85)
    body_h = Inches(3.55) if callout else Inches(3.9)

    for i, texto in enumerate(items):
        x = content_x0 + i * (col_w + gap)
        titulo_col, cuerpo_col = dividir_titulo_cuerpo(texto)
        color = PALETA_CLASIFICACION[i % len(PALETA_CLASIFICACION)]
        if not titulo_col:
            titulo_col, cuerpo_col = f"{i + 1}", texto

        cabecera = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, col_w, header_h)
        cabecera.fill.solid()
        cabecera.fill.fore_color.rgb = color
        cabecera.line.fill.background()
        cabecera.shadow.inherit = False
        cabecera.text_frame.word_wrap = True
        cabecera.text_frame.margin_left = Inches(0.12)
        cabecera.text_frame.margin_right = Inches(0.12)
        cabecera.text_frame.text = titulo_col.upper()
        cabecera.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        estilo_run(cabecera.text_frame.paragraphs[0].runs[0], 13, WHITE, negrita=True)
        cabecera.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        cuerpo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0 + header_h + Inches(0.08), col_w, body_h)
        cuerpo.fill.solid()
        cuerpo.fill.fore_color.rgb = WHITE
        cuerpo.line.color.rgb = LINE_GRAY
        cuerpo.line.width = Pt(0.75)
        cuerpo.shadow.inherit = False
        cuerpo.text_frame.word_wrap = True
        cuerpo.text_frame.margin_left = Inches(0.15)
        cuerpo.text_frame.margin_right = Inches(0.15)
        cuerpo.text_frame.margin_top = Inches(0.12)
        cuerpo.text_frame.text = cuerpo_col
        estilo_run(cuerpo.text_frame.paragraphs[0].runs[0], 12.5, CHARCOAL)

    if callout:
        anadir_callout(slide, callout, y0 + header_h + Inches(0.08) + body_h + Inches(0.15))
    anadir_pie(slide, prs, pie_izquierda, numero, MUTED_GRAY)
    return slide


def slide_comparativa(prs, titulo, bloque, pie_izquierda, numero):
    """Patrón H — dos paneles grandes enfrentados (uno oscuro, uno naranja)."""
    subtitulos, callout, items_raw = separar_bloque(bloque)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, PAPER_LIGHT)
    y0 = cabecera_contenido(slide, titulo, subtitulos[0] if subtitulos else None)

    items = [texto for tipo, texto in items_raw if tipo == "bullet"][:2]
    colores = [INK_BLACK, BRAND_ORANGE]
    gap = Inches(0.3)
    content_x0, content_w = Inches(0.75), Inches(11.83)
    panel_w = int((content_w - gap) / 2)
    panel_h = Inches(4.15) if callout else Inches(4.85)

    for i, texto in enumerate(items):
        x = content_x0 + i * (panel_w + gap)
        titulo_p, cuerpo_p = dividir_titulo_cuerpo(texto)

        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, panel_w, panel_h)
        panel.fill.solid()
        panel.fill.fore_color.rgb = colores[i]
        panel.line.fill.background()
        panel.shadow.inherit = False
        panel.text_frame.text = ""

        anadir_texto(slide, (titulo_p or f"Bloque {i + 1}").upper(), x + Inches(0.3), y0 + Inches(0.3),
                     panel_w - Inches(0.6), Inches(0.5), 16, WHITE, negrita=True)

        sub_items = [s.strip() for s in re.split(r";\s*", cuerpo_p) if s.strip()]
        if len(sub_items) > 1:
            anadir_lista(slide, sub_items, x + Inches(0.3), y0 + Inches(0.95),
                         panel_w - Inches(0.6), panel_h - Inches(1.2), 13, WHITE)
        else:
            anadir_texto(slide, cuerpo_p, x + Inches(0.3), y0 + Inches(0.95),
                         panel_w - Inches(0.6), panel_h - Inches(1.2), 13, WHITE)

    if callout:
        anadir_callout(slide, callout, y0 + panel_h + Inches(0.2))
    anadir_pie(slide, prs, pie_izquierda, numero, MUTED_GRAY)
    return slide


def slide_contenido(prs, titulo, bloque, pie_izquierda, numero):
    subtitulos, callout, items_raw = separar_bloque(bloque)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_solido(slide, PAPER_LIGHT)
    content_y0 = cabecera_contenido(slide, titulo, subtitulos[0] if subtitulos else None)

    subs = [texto for tipo, texto in items_raw if tipo == "sub"]
    tarjetas = [texto for tipo, texto in items_raw if tipo == "bullet"]

    total_palabras = sum(len(t.split()) for t in subs + tarjetas)
    if not (15 <= total_palabras <= 55):
        avisos.append(
            f"Diapositiva \"{titulo}\" tiene {total_palabras} palabras de cuerpo (rango recomendado: 15-55, sección 10.1 de la guía)."
        )

    content_x0, content_w = Inches(0.75), Inches(11.83)
    gap = Inches(0.22)

    # Subtítulos (###, poco frecuentes): línea completa encima de la cuadrícula.
    for i, sub in enumerate(subs):
        anadir_texto(slide, sub, content_x0, content_y0 + i * Inches(0.42), content_w, Inches(0.4), 18, CHARCOAL, negrita=True)
    if subs:
        content_y0 += Inches(0.42) * len(subs) + Inches(0.12)

    filas = filas_de_tarjetas(len(tarjetas)) if tarjetas else []
    max_cols = max(filas) if filas else 1
    card_w = (content_w - (max_cols - 1) * gap) / max_cols
    if len(filas) == 1:
        row_h = Inches(2.3) if callout else Inches(2.6)
    else:
        row_h = Inches(1.95) if callout else Inches(2.25)

    idx = 0
    for fila_i, n_en_fila in enumerate(filas):
        fila_w = n_en_fila * card_w + (n_en_fila - 1) * gap
        x_offset = content_x0 + int((content_w - fila_w) / 2)
        y = content_y0 + fila_i * (row_h + gap)
        for col_i in range(n_en_fila):
            texto = tarjetas[idx]
            idx += 1
            x = x_offset + col_i * (card_w + gap)

            tarjeta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, row_h)
            tarjeta.fill.solid()
            tarjeta.fill.fore_color.rgb = WHITE
            tarjeta.line.color.rgb = LINE_GRAY
            tarjeta.line.width = Pt(0.75)
            tarjeta.shadow.inherit = False
            tarjeta.text_frame.margin_left = Inches(0.22)
            tarjeta.text_frame.margin_right = Inches(0.22)
            tarjeta.text_frame.margin_top = Inches(0.1)
            tarjeta.text_frame.margin_bottom = Inches(0.1)
            tarjeta.text_frame.word_wrap = True
            tarjeta.text_frame.text = ""

            circulo = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.16), Inches(0.4), Inches(0.4))
            circulo.fill.solid()
            circulo.fill.fore_color.rgb = BRAND_ORANGE
            circulo.line.fill.background()
            circulo.shadow.inherit = False
            tf_c = circulo.text_frame
            tf_c.text = str(idx)
            tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
            estilo_run(tf_c.paragraphs[0].runs[0], 14, WHITE, negrita=True)

            titulo_tarjeta, cuerpo_tarjeta = dividir_titulo_cuerpo(texto)
            texto_y = y + Inches(0.68)
            texto_h = row_h - Inches(0.8)
            if titulo_tarjeta:
                anadir_texto(slide, titulo_tarjeta, x + Inches(0.2), texto_y, card_w - Inches(0.4), Inches(0.4),
                             15, BRAND_ORANGE, negrita=True)
                anadir_texto(slide, cuerpo_tarjeta, x + Inches(0.2), texto_y + Inches(0.4), card_w - Inches(0.4),
                             texto_h - Inches(0.4), 12.5, CHARCOAL)
            else:
                anadir_texto(slide, cuerpo_tarjeta, x + Inches(0.2), texto_y, card_w - Inches(0.4), texto_h, 13.5, CHARCOAL)

    if callout:
        bottom = content_y0 + len(filas) * (row_h + gap) if filas else content_y0 + Inches(0.3)
        anadir_callout(slide, callout, bottom)
    anadir_pie(slide, prs, pie_izquierda, numero, MUTED_GRAY)
    return slide


# --- Construcción -------------------------------------------------------------

def construir_pptx(secciones, titulo_presentacion, subtitulo, etiqueta, ruta_salida, pie_izquierda):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    limpias = [extraer_patron(t) + (i,) for t, i in secciones]  # (titulo_limpio, patron_tag, items)
    clasificadas = [(t, tag, i, clasificar_seccion(t, i)) for t, tag, i in limpias]
    total_bloques = sum(1 for _, _, _, c in clasificadas if c == "separador")

    RENDERERS_CONTENIDO = {
        "recorrido": slide_recorrido,
        "clasificacion": slide_clasificacion,
        "comparativa": slide_comparativa,
        "reglaexcepciones": slide_reglaexcepciones,
    }

    numero = 1
    slide_portada(prs, titulo_presentacion, subtitulo, etiqueta, pie_izquierda, numero)

    bloque_actual = 0
    for titulo, patron_tag, items, clase in clasificadas:
        if clase == "separador":
            bloque_actual += 1
            numero += 1
            slide_separador(prs, titulo, pie_izquierda, numero, bloque_actual, total_bloques)
            continue

        if clase == "protocolo":
            numero += 1
            slide_protocolo(prs, titulo, items, pie_izquierda, numero)
            continue

        if clase == "cierre":
            numero += 1
            slide_cierre(prs, quitar_prefijo_separador(titulo) or titulo, items, pie_izquierda, numero)
            continue

        renderer = RENDERERS_CONTENIDO.get(patron_tag, slide_contenido)
        # El subtítulo (>) y el callout (- !) no son viñetas de contenido:
        # se trocean aparte y se reenganchan siempre a la primera/última
        # diapositiva, para que un callout nunca acabe solo en una
        # diapositiva fantasma "(2/2)".
        subtitulo_items = [it for it in items if it[0] == "subtitulo"]
        callout_items = [it for it in items if it[0] == "callout"]
        contenido_items = [it for it in items if it[0] not in ("subtitulo", "callout")]
        bloques = [contenido_items[i:i + MAX_VINETAS_POR_DIAPOSITIVA] for i in range(0, len(contenido_items), MAX_VINETAS_POR_DIAPOSITIVA)] or [[]]
        for idx, bloque in enumerate(bloques):
            numero += 1
            sufijo = f" ({idx + 1}/{len(bloques)})" if len(bloques) > 1 else ""
            bloque_final = (subtitulo_items if idx == 0 else []) + bloque + (callout_items if idx == len(bloques) - 1 else [])
            renderer(prs, titulo + sufijo, bloque_final, pie_izquierda, numero)

    prs.save(ruta_salida)
    return len(prs.slides)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True)
    ap.add_argument("--output", required=True)
    ap.add_argument("--titulo", required=True, help="Título de portada (máx. 12 palabras recomendadas)")
    ap.add_argument("--subtitulo", default="")
    ap.add_argument(
        "--etiqueta",
        default="",
        help="Etiqueta corta sobre el título de portada (p. ej. 'VCF · TSEAS' para VCF, 'MET · TSEAS' para MET)",
    )
    ap.add_argument(
        "--pie",
        default="FORMDEPOR",
        help="Texto del pie de página (footer_left) (p. ej. 'VCF · TSEAS · FORMDEPOR' para VCF, 'MET · TSEAS · FORMDEPOR' para MET)",
    )
    args = ap.parse_args()

    with open(args.input, encoding="utf-8") as f:
        md_texto = f.read()

    secciones = parsear_secciones(md_texto)
    if not secciones:
        raise SystemExit(f"No se encontraron secciones '##' en {args.input} — nada que convertir.")

    n = construir_pptx(secciones, args.titulo, args.subtitulo, args.etiqueta, args.output, args.pie)

    print(f"OK: {args.output} generado con {len(secciones)} secciones de origen ({n} diapositivas), sistema visual UD00 style.")
    if avisos:
        print("\nAvisos de estilo (revisar antes de dar la presentación por buena):")
        for aviso in avisos:
            print(f"  - {aviso}")


if __name__ == "__main__":
    main()
